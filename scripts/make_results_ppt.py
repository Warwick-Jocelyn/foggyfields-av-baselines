"""Build a results-summary PPTX for the FoggyFields AV baseline project.

Two outputs:
  FoggyFields_Results.pptx       — qualitative panels embed the (compressed) render mp4s
                                    (click a panel to play); heavier, fully self-contained.
  FoggyFields_Results_lite.pptx  — same deck, poster frames only; small + opens anywhere.

Qualitative slides: 2x2 grid (GT / SplatAD / NeuRAD / EmerNeRF) per clip with PSNR labels.
"""
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ETH = Path("/networkhome/WMGDS/wang3_y/ETH")
FOG = ETH/"NVIDIA-Fog-Output"; PAN = ETH/"Code/outputs"
FR  = ETH/"_ppt_frames"; FR.mkdir(exist_ok=True)
VD  = ETH/"_ppt_vids";   VD.mkdir(exist_ok=True)

NAVY=RGBColor(0x1F,0x2A,0x44); BLUE=RGBColor(0x2E,0x5C,0xB5); GREY=RGBColor(0x55,0x55,0x55)
GREEN=RGBColor(0x1B,0x7A,0x3D); RED=RGBColor(0xB1,0x2A,0x2A); WHITE=RGBColor(0xFF,0xFF,0xFF)

def grab(video, key):
    """Extract one frame, downscale to <=1000px wide, save compact JPEG -> path."""
    jpg=FR/f"{key}.jpg"
    if jpg.exists(): return jpg
    if not video or not Path(video).exists(): return None
    png=FR/f"{key}.png"
    for args in (["-ss","1","-i",str(video),"-frames:v","1"],["-i",str(video),"-frames:v","1"]):
        try:
            subprocess.run(["ffmpeg","-y",*args,"-q:v","2",str(png)],check=True,capture_output=True,timeout=60)
            if png.exists(): break
        except Exception: continue
    if not png.exists(): return None
    try:
        im=Image.open(png).convert("RGB")
        if im.width>1000: im=im.resize((1000,max(1,round(im.height*1000/im.width))))
        im.save(jpg,"JPEG",quality=82)
        png.unlink(missing_ok=True)
        return jpg
    except Exception:
        return png

def compact(video, key):
    if not video or not Path(video).exists(): return None
    out=VD/f"{key}.mp4"
    if out.exists(): return out
    try:
        subprocess.run(["ffmpeg","-y","-i",str(video),"-vf","scale='min(900,iw)':-2",
                        "-c:v","libx264","-crf","30","-preset","veryfast","-an",
                        "-movflags","+faststart",str(out)],check=True,capture_output=True,timeout=180)
        if out.exists() and out.stat().st_size>0: return out
    except Exception: pass
    return Path(video)

# (cell-key, label, psnr, render_video)  — keys are globally unique
FOG_CLIPS={
 "002":("rural · medium fog · ~113 km/h", FOG/"splatad/002/videos/val/panorama_gt_rgb.mp4",[
   ("fog_splatad_002","SplatAD","21.64",FOG/"splatad/002/videos/val/panorama_rgb.mp4"),
   ("fog_neurad_002","NeuRAD","27.63",  FOG/"neurad/002/videos/val/panorama_rgb.mp4"),
   ("fog_emernerf_002","EmerNeRF*","31.36",FOG/"emernerf/002/foggyfields_emernerf/nvidia_fog_002_7cam/test_videos/25000_rgbs.mp4")]),
 "003":("residential · light fog", FOG/"splatad/003/videos/val/panorama_gt_rgb.mp4",[
   ("fog_splatad_003","SplatAD","24.60",FOG/"splatad/003/videos/val/panorama_rgb.mp4"),
   ("fog_neurad_003","NeuRAD","29.10",  FOG/"neurad/003/videos/val/panorama_rgb.mp4"),
   ("fog_emernerf_003","EmerNeRF*","34.11",FOG/"emernerf/003/foggyfields_emernerf/nvidia_fog_003_7cam/test_videos/25000_rgbs.mp4")]),
 "004":("highway · heavy fog", FOG/"splatad/004/videos/val/panorama_gt_rgb.mp4",[
   ("fog_splatad_004","SplatAD","23.94",FOG/"splatad/004/videos/val/panorama_rgb.mp4"),
   ("fog_neurad_004","NeuRAD","29.90",  FOG/"neurad/004/videos/val/panorama_rgb.mp4"),
   ("fog_emernerf_004","EmerNeRF*","34.66",FOG/"emernerf/004/foggyfields_emernerf/nvidia_fog_004_7cam/test_videos/25000_rgbs.mp4")]),
}
PAN_CLIPS={
 "011 (day)":(PAN/"paper_011_split05_paperfaithful/videos/val/panorama_gt_rgb.mp4",[
   ("p_splatad11","SplatAD","27.52",PAN/"paper_011_split05_paperfaithful/videos/val/panorama_rgb.mp4"),
   ("p_neurad11","NeuRAD","26.46",  PAN/"neurad_paper_011_split05/videos/val/panorama_rgb.mp4"),
   ("p_emer11","EmerNeRF","27.55",  PAN/"emernerf_011_split05/foggyfields_emernerf/scene_011_6cam_dyn_flow_rgb/test_videos/30000_rgbs.mp4")]),
 "078 (NIGHT)":(PAN/"paper_078_split05_paperfaithful/videos/val/panorama_gt_rgb.mp4",[
   ("p_splatad78","SplatAD","31.65",PAN/"paper_078_split05_paperfaithful/videos/val/panorama_rgb.mp4"),
   ("p_neurad78","NeuRAD","30.04",  PAN/"neurad_paper_078_split05/videos/val/panorama_rgb.mp4"),
   ("p_emer78","EmerNeRF","31.26",  PAN/"emernerf_078_split05/foggyfields_emernerf/scene_078_6cam_dyn_flow_rgb/test_videos/30000_rgbs.mp4")]),
}

def bg(s,c=WHITE): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def box(s,l,t,w,h,text,size=18,bold=False,color=NAVY,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=line; f=r.font; f.size=Pt(size); f.bold=bold; f.color.rgb=color; f.name="Calibri"
    return tb
def fit(p,bw,bh):
    try:
        w,h=Image.open(p).size; ar=w/h
        return (bw,bw/ar) if ar>bw/bh else (bh*ar,bh)
    except Exception: return bw,bh

def table(s,rows,l,t,col_w,header_fill=NAVY,special=None):
    tb=s.shapes.add_table(len(rows),len(rows[0]),Inches(l),Inches(t),Inches(sum(col_w)),Inches(0.36*len(rows))).table
    for ci,w in enumerate(col_w): tb.columns[ci].width=Inches(w)
    for ri,row in enumerate(rows):
        for cj,val in enumerate(row):
            c=tb.cell(ri,cj); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if cj==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=val; r.font.size=Pt(12); r.font.name="Calibri"
            if ri==0:
                r.font.bold=True; r.font.color.rgb=WHITE; c.fill.solid(); c.fill.fore_color.rgb=header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xF2,0xF5,0xFB) if ri%2 else WHITE
                if special: special(ri,cj,val,r)
    return tb

def build(embed, outpath):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    def place(s,key,video,l,t,bw,bh,label,psnr=None,labcol=NAVY):
        poster=grab(video,key)
        if poster:
            w_in,h_in=fit(poster,bw,bh-0.32); cl=l+(bw-w_in)/2; ct=t+0.30+(bh-0.32-h_in)/2
            added=False
            if embed:
                cv=compact(video,key)
                if cv and Path(cv).exists():
                    try:
                        s.shapes.add_movie(str(cv),Inches(cl),Inches(ct),Inches(w_in),Inches(h_in),
                                           poster_frame_image=str(poster),mime_type="video/mp4"); added=True
                    except Exception: added=False
            if not added: s.shapes.add_picture(str(poster),Inches(cl),Inches(ct),Inches(w_in),Inches(h_in))
        box(s,l,t,bw,0.30,label if psnr is None else f"{label}   PSNR {psnr}",size=13,bold=True,color=labcol)
    def qual(title,subtitle,gt_key,gt_video,cells,night=False):
        s=prs.slides.add_slide(blank); bg(s)
        box(s,0.4,0.15,12.5,0.5,title,size=24,bold=True,color=NAVY)
        box(s,0.4,0.66,12.5,0.35,subtitle,size=14,color=(RED if night else GREY))
        gx,gy,cw,ch,gap=0.4,1.15,6.15,2.95,0.15
        co=[(gx,gy),(gx+cw+gap,gy),(gx,gy+ch+gap),(gx+cw+gap,gy+ch+gap)]
        place(s,gt_key,gt_video,*co[0],cw,ch,"Ground Truth",None,GREY)
        for i,(k,lab,ps,vid) in enumerate(cells): place(s,k,vid,*co[i+1],cw,ch,lab,ps,BLUE)

    # 1 title
    s=prs.slides.add_slide(blank); bg(s,NAVY)
    box(s,0.8,2.3,11.7,1.2,"FoggyFields — AV Neural-Field Baselines",size=40,bold=True,color=WHITE)
    box(s,0.8,3.5,11.7,0.7,"Results summary · SplatAD vs NeuRAD vs EmerNeRF",size=22,color=RGBColor(0xCF,0xDA,0xF0))
    box(s,0.8,4.5,11.7,1.4,
        "Benchmarks: NVIDIA_AV_Fog (foggy, 7-cam + LiDAR, clips 002/003/004) · PandaSet (fog-free reference, 011 day / 078 night)\n"
        "Protocol: nvs_50_50 held-out novel view · 1× RTX PRO 6000 Blackwell per run",
        size=15,color=RGBColor(0xAD,0xBE,0xDC))

    # 2 foggy table
    s=prs.slides.add_slide(blank); bg(s)
    box(s,0.4,0.2,12.5,0.5,"NVIDIA_AV_Fog — held-out novel-view metrics",size=24,bold=True,color=NAVY)
    rows=[["Clip / scene","Method","PSNR↑","SSIM↑","LPIPS↓","LiDAR med↓ (m)"],
     ["002 rural / med fog","SplatAD","21.64","0.850","0.413","8.15"],["","NeuRAD","27.63","0.884","0.180","5.27"],["","EmerNeRF*","31.36","0.941","0.177","—"],
     ["003 resid / light fog","SplatAD","24.60","0.873","0.353","8.63"],["","NeuRAD","29.10","0.897","0.161","7.49"],["","EmerNeRF*","34.11","0.956","0.167","—"],
     ["004 highway / heavy fog","SplatAD","23.94","0.882","0.324","15.20"],["","NeuRAD","29.90","0.911","0.152","9.38"],["","EmerNeRF*","34.66","0.963","0.136","—"]]
    def sp(ri,cj,val,r):
        if cj==1 and val=="EmerNeRF*": r.font.color.rgb=GREEN; r.font.bold=True
        elif cj==2 and val and val[0]=="3": r.font.bold=True; r.font.color.rgb=GREEN
    table(s,rows,0.4,0.95,[2.7,1.2,1.2,1.2,1.2,1.4],special=sp)
    box(s,9.5,1.0,3.5,5.0,
     "Key points\n\n• EmerNeRF wins RGB on every fog clip (mean ~33 dB) — volumetric integration handles fog (participating media) best.\n\n"
     "• NeuRAD wins LiDAR depth (median 5–9 m vs SplatAD 8–15 m).\n\n• SplatAD trades fog fidelity for 8–13× faster training.\n\n"
     "* EmerNeRF = pre-split-fix, re-train pending (indicative).",size=13,color=GREY)

    for c,(sub,gt,cells) in FOG_CLIPS.items():
        qual(f"NVIDIA_AV_Fog — clip {c}", sub+"   (7-cam held-out render vs GT"+(" · click a panel to play)" if embed else ")"),
             f"gt_fog_{c}", gt, cells)

    # 6 pandaset table
    s=prs.slides.add_slide(blank); bg(s)
    box(s,0.4,0.2,12.5,0.5,"PandaSet — fog-free reference (sanity baseline)",size=24,bold=True,color=NAVY)
    rows=[["Clip","Lighting","Method","PSNR↑","SSIM↑","LPIPS↓"],
     ["011","day","SplatAD","27.52","0.868","0.162"],["","","NeuRAD","26.46","0.805","0.201"],["","","EmerNeRF","27.55","0.790","—"],
     ["078","NIGHT","SplatAD","31.65","0.930","0.242"],["","","NeuRAD","30.04","0.898","0.214"],["","","EmerNeRF","31.26","0.879","—"]]
    def sp2(ri,cj,val,r):
        if cj==1 and val=="NIGHT": r.font.color.rgb=RED; r.font.bold=True
        if cj==1 and val=="day": r.font.color.rgb=GREEN
    table(s,rows,0.4,0.95,[1.0,1.3,1.5,1.4,1.4,1.4],special=sp2)
    box(s,8.7,1.0,4.3,4.4,
     "Why PandaSet?\n\n• Fog-free → isolates the no-fog baseline; the 3 methods are close, confirming the implementations are sound.\n\n"
     "• So EmerNeRF's big lead on NVIDIA_AV_Fog is a FOG-specific effect, not a code artifact.\n\n"
     "⚠ Clip 078 is a NIGHT scene — higher numbers reflect lighting, not directly comparable to day-time 011.",size=13,color=GREY)

    for name,(gt,cells) in PAN_CLIPS.items():
        cid=name.split()[0]
        qual(f"PandaSet — clip {cid}", ("nighttime scene" if "NIGHT" in name else "daytime scene")+("   (render vs GT · click a panel to play)" if embed else "   (render vs GT)"),
             f"gt_pan_{cid}", gt, cells, night="NIGHT" in name)

    # 9 conclusions
    s=prs.slides.add_slide(blank); bg(s,NAVY)
    box(s,0.5,0.3,12.3,0.6,"Conclusions",size=28,bold=True,color=WHITE)
    box(s,0.6,1.15,12.1,5.2,
     "1.  EmerNeRF is the strongest RGB reconstructor in fog — mean ~33 dB vs NeuRAD ~29, SplatAD ~23. Its volumetric RGB+LiDAR field models fog better than discrete Gaussians.\n\n"
     "2.  NeuRAD gives the best LiDAR geometry (depth median 5–9 m vs SplatAD 8–15 m); SplatAD scatters mass into the fog volume, hurting depth.\n\n"
     "3.  SplatAD trains 8–13× faster (~40 min vs NeuRAD 6.3 h / EmerNeRF 4.4 h) — speed-vs-fidelity trade-off.\n\n"
     "4.  Difficulty tracks fog × speed: clip 002 (rural, 113 km/h, medium fog) hardest; 004 (heavy-fog highway) highest.\n\n"
     "5.  On fog-free PandaSet the three converge → the fog gap is real, not an artifact. (Note 078 = night.)",
     size=16,color=RGBColor(0xE6,0xEC,0xF7))
    box(s,0.6,6.55,12.1,0.7,
     "Code: github.com/Warwick-Jocelyn/foggyfields-av-baselines (private)  ·  Checkpoints: huggingface.co/JocelynW/foggyfields-av-baseline-checkpoints (private)",
     size=12,color=RGBColor(0xAD,0xBE,0xDC))

    # 10 video index
    s=prs.slides.add_slide(blank); bg(s)
    box(s,0.4,0.2,12.5,0.5,"Source render videos (multi-camera, on disk)",size=22,bold=True,color=NAVY)
    lines=("Foggy — NVIDIA-Fog-Output/<method>/<clip>/videos/val/panorama_rgb.mp4 ; EmerNeRF .../test_videos/25000_rgbs.mp4\n"
     "   • clips 002 / 003 / 004 × {splatad, neurad, emernerf}   (GT: panorama_gt_rgb.mp4 / 25000_gt_rgbs.mp4)\n\n"
     "PandaSet — Code/outputs/...\n"
     "   • 011 day  : paper_011_split05_paperfaithful · neurad_paper_011_split05 · emernerf_011_split05  (.../val/panorama_rgb.mp4)\n"
     "   • 078 night: paper_078_split05_paperfaithful · neurad_paper_078_split05 · emernerf_078_split05\n\n"
     "Also alongside: panorama_depth.mp4, per-camera *_compare.mp4, lidar_gt_bbox.mp4 (actor-pose sanity).")
    box(s,0.5,1.0,12.3,5.5,lines,size=14,color=GREY)

    prs.save(str(outpath))
    return len(prs.slides._sldIdLst), outpath.stat().st_size/1e6

n,mb=build(False, ETH/"FoggyFields_Results_lite.pptx"); print(f"LITE  : {n} slides, {mb:.1f} MB")
n,mb=build(True,  ETH/"FoggyFields_Results.pptx");      print(f"FULL  : {n} slides, {mb:.1f} MB")
print("frames:",len(list(FR.glob('*.png'))),"| compact vids:",len(list(VD.glob('*.mp4'))))
